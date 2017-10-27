from pptx import Presentation
from openpyxl import load_workbook
from openpyxl.utils.dataframe import dataframe_to_rows


class PPT:
    def __init__(self, input_file='template.pptx'):
        self.prs = Presentation(input_file)
        self.slides = self.prs.slides

    def analyze_layouts(self, outputfile='output.pptx'):
        # 遍历每个版式与占位符
        for s, layout in enumerate(self.prs.slide_layouts):
            slide = self.prs.slides.add_slide(layout)

            # 是否有标题占位符
            try:
                title = slide.shapes.title
                title.text = f'{s}样式-标题'
            except AttributeError:
                print(f'>>> page {i} has no title')

            # 将其他占位符(placeholders)命名为x样式x号
            for each in slide.placeholders:
                each.text = f'{s}样式-{each.placeholder_format.idx}号'

        # 保存
        self.save(outputfile)

    def analyze_slides(self, output_file='output.pptx'):
        # 遍历每页与每个shape
        for p, slide in enumerate(self.slides):
            for i, shape in enumerate(slide.shapes):
                shape.text = f'{p}页-{i}号'

        # 保存
        self.save(output_file)

    def df2table(self, df, page_idx, shape_idx):
        """将df中数据填入指定表格(通过指定页与shape)
        :param df: DataFrame对象
        :param page_idx: 表格所在页数编号
        :param shape_idx: 表格所在shape编号
        """
        tb = self.slides[page_idx].shapes[shape_idx].table

        # 确定表格行、列数
        rows, cols = df.shape

        # 填写表头
        columns = list(df.columns)
        for col, value in enumerate(columns):
            cell = tb.cell(0, col)
            cell.text = value

        # 填写数据
        df_matrix = df.as_matrix()
        for row in range(rows):
            for col in range(cols):
                value = df_matrix[row, col]
                cell = tb.cell(row + 1, col)
                cell.text = str(value)

    def text(self, txt, page_idx, shape_idx):
        """将文字填入指定shape
        :param txt: 需填入的文字
        :param page_idx: shape所在页数编号
        :param shape_idx: shape编号
        """
        self.slides[page_idx].shapes[shape_idx].text = txt

    def save(self, output_file='output.pptx'):
        self.prs.save(output_file)


class Excel:
    def __init__(self, input_file='template.xlsx'):
        self.wb = load_workbook(input_file)

    def df2sheet(self, df, sheet_name):
        """将df贴到指定sheet上
        :param df: DataFrame对象
        :param sheet_name:sheet名称
        """
        ws = self.wb.get_sheet_by_name(sheet_name)
        for row in dataframe_to_rows(df, header=False):
            ws.append(row)

    def save(self, output_name='data.xlsx'):
        self.wb.save(output_name)


if __name__ == '__main__':
    path = r'E:\dell\Documents\Python_Projects\report_factory_3'
    ppt = PPT(f'{path}/template.pptx')
    ppt.analyze_slides()
